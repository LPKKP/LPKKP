"""Generate Pembentangan KKP PowerPoint for management."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colours — professional government tone
NAVY = RGBColor(0x1A, 0x3A, 0x5C)
BLUE = RGBColor(0x2E, 0x6B, 0xA4)
TEAL = RGBColor(0x0D, 0x5C, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
ACCENT_GOLD = RGBColor(0xC9, 0xA2, 0x27)

OUTPUT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "Pembentangan_Sistem_Laporan_KKP.pptx",
)


def set_run_font(run, size=18, bold=False, color=DARK, name="Calibri"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_bar(slide, title_text):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_run_font(run, size=28, bold=True, color=WHITE)


def add_bullets(slide, items, left=0.7, top=1.45, width=8.8, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = item[1] if isinstance(item, tuple) else 0
        text = item[0] if isinstance(item, tuple) else item
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = text
        indent = 0 if p.level == 0 else 1
        set_run_font(run, size=size - indent * 2, bold=False, color=DARK)
    return box


def add_footer(slide, text="Sistem Laporan Prestasi KKP — Sulit Dalaman"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(9), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run_font(run, size=10, color=GRAY)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.8), Inches(10), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GOLD
    accent.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "SISTEM LAPORAN PRESTASI"
    set_run_font(r, size=36, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "KESELAMATAN & KESIHATAN PEKERJAAN (KKP)"
    set_run_font(r2, size=32, bold=True, color=WHITE)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.4), Inches(1.2))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "Pembentangan kepada Pengurusan Atasan"
    set_run_font(sr, size=20, color=RGBColor(0xCC, 0xDD, 0xEE))
    sp2 = stf.add_paragraph()
    sp2.alignment = PP_ALIGN.CENTER
    sr2 = sp2.add_run()
    sr2.text = "Inisiatif Digitalisasi Pelaporan Suku Tahunan"
    set_run_font(sr2, size=16, color=RGBColor(0xAA, 0xBB, 0xCC))


def slide_content(prs, title, bullets, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    add_bullets(slide, bullets)
    add_footer(slide, footer or "Sistem Laporan Prestasi KKP")


def slide_two_column(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)

    for x, heading, items in [
        (0.5, left_title, left_items),
        (5.2, right_title, right_items),
    ]:
        hbox = slide.shapes.add_textbox(Inches(x), Inches(1.35), Inches(4.3), Inches(0.5))
        hp = hbox.text_frame.paragraphs[0]
        hr = hp.add_run()
        hr.text = heading
        set_run_font(hr, size=16, bold=True, color=TEAL)
        add_bullets(slide, items, left=x, top=1.85, width=4.3, height=4.8, size=15)

    add_footer(slide)


def slide_cost_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Anggaran Kos Projek")

    intro = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.6))
    ip = intro.text_frame.paragraphs[0]
    ir = ip.add_run()
    ir.text = "Anggaran berdasarkan skala penggunaan semasa dan pasaran Sabah / Malaysia (2026)"
    set_run_font(ir, size=14, color=GRAY)

    rows, cols = 5, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.95), Inches(8.8), Inches(2.2))
    table = table_shape.table
    headers = ["Komponen", "Kos Anggaran (RM)", "Catatan"]
    data = [
        ["Pembangunan (dalaman)", "0 – 15,000", "Kos peluang masa kakitangan; tiada invois langsung"],
        ["Pembangunan (outsource Sabah)", "15,000 – 25,000", "Nilai pasaran jika dibangunkan oleh vendor"],
        ["Operasi tahunan (Firebase / hosting)", "0 – 500 / tahun", "Percuma pada skala semasa"],
        ["Penyelenggaraan tahunan", "0 – 5,000 / tahun", "Sokongan teknikal & kemas kini kecil"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run_font(r, size=12, bold=True, color=WHITE)

    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_font(run, size=11, color=DARK)

    highlight = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.35), Inches(8.8), Inches(1.1)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    highlight.line.color.rgb = TEAL
    htf = highlight.text_frame
    htf.word_wrap = True
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run()
    hr.text = (
        "Ringkasan: Kos operasi tahunan hampir sifar (RM 0–500). "
        "Berbanding perisian komersial KKP (RM 20,000–100,000+/tahun), "
        "sistem ini memberi nilai tinggi pada kos rendah."
    )
    set_run_font(hr, size=14, bold=True, color=TEAL)

    add_footer(slide)


def slide_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Perbandingan Pendekatan")

    rows, cols = 5, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(8.8), Inches(3.5))
    table = table_shape.table
    headers = ["Pendekatan", "Kos Tahunan", "Keterangan"]
    data = [
        ["Sistem ini (Laporan KKP)", "RM 0 – 3,000", "Digital, terpusat, jejak audit, PDF automatik"],
        ["Excel + e-mel", "RM 0", "Manual, sukar dijejak, tiada kawalan berpusat"],
        ["Perisian KKP komersial", "RM 20,000 – 200,000+", "Mahal, mungkin melebihi keperluan"],
        ["Sistem enterprise (tender)", "RM 100,000 – 500,000+", "Jangka masa panjang, kos tinggi"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run_font(r, size=12, bold=True, color=WHITE)

    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_font(run, size=11, color=DARK)

    add_footer(slide)


def slide_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Aliran Kerja Sistem")

    steps = [
        ("1", "Daftar", "Bahagian/Zon/Unit mendaftar akaun"),
        ("2", "Kelulusan", "Pentadbir meluluskan pengguna baharu"),
        ("3", "Log Masuk", "Pengguna log masuk (e-mel atau bahagian)"),
        ("4", "Isi Borang", "Laporan suku KKP diisi dalam talian"),
        ("5", "PDF Automatik", "Laporan PDF dijana & dimuat turun"),
        ("6", "Pemantauan", "Pengurusan pantau status penghantaran"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 3.15
        y = 1.55 + row * 2.55

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.55), Inches(0.55)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        set_run_font(cr, size=16, bold=True, color=WHITE)

        tbox = slide.shapes.add_textbox(Inches(x + 0.65), Inches(y), Inches(2.35), Inches(1.0))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = title
        set_run_font(tr, size=15, bold=True, color=NAVY)
        tp2 = ttf.add_paragraph()
        tr2 = tp2.add_run()
        tr2.text = desc
        set_run_font(tr2, size=12, color=GRAY)

    add_footer(slide)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(3.5))
    tf = tbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Kesimpulan"
    set_run_font(r, size=32, bold=True, color=WHITE)

    points = [
        "Sistem ini memudahkan pematuhan pelaporan KKP suku tahunan",
        "Memberi pandangan berpusat kepada pihak pengurusan",
        "Kos operasi rendah dan penyelenggaraan ringkas",
        "Sedia untuk fasa percubaan / penggunaan rasmi",
    ]
    for pt in points:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        r2 = p2.add_run()
        r2.text = "•  " + pt
        set_run_font(r2, size=18, color=RGBColor(0xCC, 0xDD, 0xEE))

    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.8))
    tfp = thanks.text_frame.paragraphs[0]
    tfp.alignment = PP_ALIGN.CENTER
    tr = tfp.add_run()
    tr.text = "Terima Kasih  |  Saya bersedia menjawab sebarang pertanyaan"
    set_run_font(tr, size=18, bold=True, color=ACCENT_GOLD)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)

    slide_content(
        prs,
        "Latar Belakang & Cabaran",
        [
            "Pelaporan prestasi KKP suku tahunan masih bergantung pada kaedah manual atau e-mel",
            "Sukar untuk menjejak status penghantaran merentas Bahagian / Zon / Unit",
            "Tiada pandangan berpusat bagi pihak pengurusan",
            "Data dan dokumen tidak terstandard dan sukar diaudit",
            "Proses penyediaan PDF memakan masa dan berulang",
        ],
    )

    slide_content(
        prs,
        "Objektif Sistem",
        [
            "Mendigitalkan pelaporan suku Keselamatan & Kesihatan Pekerjaan (KKP)",
            "Menyediakan borang standard untuk semua bahagian",
            "Menjana laporan PDF secara automatik selepas penghantaran",
            "Membolehkan pihak pengurusan memantau status penghantaran",
            "Meningkatkan akauntabiliti melalui kawalan akses dan jejak audit",
        ],
    )

    slide_content(
        prs,
        "Penyelesaian Yang Dibangunkan",
        [
            "Platform web — boleh diakses melalui pelayar internet",
            "Pengguna log masuk dengan e-mel atau nama bahagian + kata laluan",
            "Borang laporan suku dengan Seksyen A hingga E (statistik, kemalangan, penyakit, aktiviti)",
            "PDF dijana secara automatik — data disimpan, PDF tidak disimpan berulang",
            "Panel pentadbir untuk kelulusan pengguna, pengurusan bahagian & jejak audit",
        ],
    )

    slide_content(
        prs,
        "Ciri-ciri Utama",
        [
            "Pendaftaran & kelulusan pengguna oleh pentadbir",
            "Borang laporan KKP suku dengan panduan pengisian",
            "Penjanaan PDF automatik selepas penghantaran",
            "Dashboard — statistik penghantaran mengikut suku",
            "Senarai data laporan — lihat, muat turun, eksport",
            "Jejak audit — rekod log masuk, penghantaran & tindakan pentadbir",
            "Tetapan semula kata laluan oleh pentadbir",
        ],
    )

    slide_workflow(prs)

    slide_two_column(
        prs,
        "Manfaat Kepada Organisasi",
        "Kecekapan",
        [
            "Borang standard untuk semua bahagian",
            "PDF dijana serta-merta selepas hantar",
            "Kemas kini laporan suku yang sama dibenarkan",
            "Kurang bergantung pada e-mel & fail berselerak",
        ],
        "Kawalan & Ketelusan",
        [
            "Dashboard status penghantaran suku",
            "Akses mengikut bahagian — pengguna lihat data sendiri sahaja",
            "Kelulusan pengguna baharu oleh pentadbir",
            "Jejak audit untuk akauntabiliti",
        ],
    )

    slide_content(
        prs,
        "Kawalan Keselamatan & Tadbir Urus",
        [
            "Log masuk diperlukan untuk semua akses sistem",
            "Peranan pengguna: Pentadbir, Pengguna, Menunggu, Ditolak",
            "Pengguna baharu mesti diluluskan sebelum boleh menggunakan sistem",
            "Data laporan dihadkan mengikut bahagian (pentadbir akses semua)",
            "Jejak audit merekod: log masuk, penghantaran laporan, pendaftaran, tindakan admin",
            "Peraturan keselamatan Firestore mengawal akses pangkalan data",
        ],
    )

    slide_cost_table(prs)
    slide_comparison(prs)

    slide_content(
        prs,
        "Status Semasa",
        [
            "Sistem telah dibangunkan dan berfungsi",
            "Modul utama siap: log masuk, pendaftaran, borang, data, admin, profil",
            "Jejak audit telah ditambah untuk pemantauan aktiviti",
            "Panduan pengisian laporan (PDF) tersedia untuk pengguna",
            "Peraturan keselamatan Firestore perlu disahkan/dikemas kini di Firebase Console",
            "Sedia untuk fasa percubaan dengan bahagian terpilih",
        ],
    )

    slide_content(
        prs,
        "Langkah Seterusnya (Cadangan)",
        [
            "Kelulusan pengurusan untuk fasa percubaan / penggunaan rasmi",
            "Penetapan pentadbir sistem untuk kelulusan pengguna",
            "Latihan ringkas untuk penyelaras bahagian",
            "Arahan rasmi: laporan suku KKP melalui sistem ini",
            "Penambahbaikan fasa seterusnya: notifikasi e-mel, analitik lanjutan",
        ],
    )

    slide_content(
        prs,
        "Demonstrasi Sistem (Live Demo)",
        [
            "1.  Pendaftaran pengguna baharu → kelulusan pentadbir",
            "2.  Log masuk sebagai pengguna bahagian",
            "3.  Isi & hantar borang laporan suku KKP",
            "4.  Muat turun PDF automatik",
            "5.  Dashboard — lihat statistik penghantaran",
            "6.  Halaman Data — senarai & tapis laporan",
            "7.  Panel Admin — pengurusan pengguna & jejak audit",
        ],
    )

    slide_closing(prs)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    build()
